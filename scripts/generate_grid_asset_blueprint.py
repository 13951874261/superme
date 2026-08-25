from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

OUT = "artifacts/2030电网资产管理蓝图-国网风格v2.pptx"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

WHITE = RGBColor(255, 255, 255)
BG = RGBColor(245, 250, 249)
GREEN = RGBColor(0, 118, 105)
GREEN_DARK = RGBColor(0, 86, 78)
GREEN_MID = RGBColor(36, 151, 137)
GREEN_LIGHT = RGBColor(224, 243, 239)
BLUE = RGBColor(0, 105, 171)
BLUE_DARK = RGBColor(0, 72, 125)
BLUE_LIGHT = RGBColor(228, 242, 249)
CYAN = RGBColor(44, 171, 183)
ORANGE = RGBColor(224, 117, 22)
GRAY = RGBColor(79, 96, 101)
GRAY_LIGHT = RGBColor(199, 218, 217)

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG


def add_shape(x, y, w, h, fill, line=GRAY_LIGHT, kind=MSO_SHAPE.RECTANGLE, width=0.8):
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(width)
    return s


def add_text(x, y, w, h, value, size=10, color=GRAY, bold=False, align=PP_ALIGN.LEFT, margin=0.03):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = s.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = value
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return s


def add_line(x1, y1, x2, y2, color=GRAY_LIGHT, width=0.8):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    s.line.color.rgb = color
    s.line.width = Pt(width)
    return s


def dot(x, y, color=GREEN):
    return add_shape(x, y, 0.09, 0.09, color, color, MSO_SHAPE.OVAL)


# Header
add_shape(0, 0, 13.333, 0.72, WHITE, WHITE)
add_line(0.42, 0.7, 12.9, 0.7, GREEN, 1.1)
add_shape(0.55, 0.14, 0.42, 0.42, GREEN_LIGHT, GREEN, MSO_SHAPE.OVAL, 1.0)
add_text(0.55, 0.22, 0.42, 0.18, "电网", 7.5, GREEN_DARK, True, PP_ALIGN.CENTER)
add_text(1.12, 0.1, 10.8, 0.42, "2030电网资产管理智慧蓝图：AI原生与数字孪生新时代", 21, GREEN_DARK, True, PP_ALIGN.CENTER)
add_text(11.95, 0.22, 0.85, 0.2, "战略展望", 8, GREEN, True, PP_ALIGN.RIGHT)

# Top left autonomy block
add_shape(0.45, 0.92, 4.18, 1.25, WHITE, GREEN_MID, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0)
add_text(0.68, 1.02, 3.72, 0.32, "超级智能体自治阶段", 16, GREEN_DARK, True)
add_text(0.68, 1.42, 3.68, 0.52, "实现自感知、业财自协同、决策自优化与知识自演进，推动资产管理由被动审批迈向自主规划。", 9.2, GRAY)

# Top right capability block
add_shape(8.68, 0.92, 4.2, 1.25, WHITE, BLUE, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0)
capabilities = [
    (8.86, 1.06, "感知与协同", "自感知 · 业财自协同"),
    (10.74, 1.06, "决策与进化", "决策自优化 · 知识自演进"),
]
for x, y, title, desc in capabilities:
    add_shape(x, y, 1.78, 0.44, BLUE_LIGHT, BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(x, y + 0.02, 1.78, 0.18, title, 9, BLUE_DARK, True, PP_ALIGN.CENTER)
    add_text(x, y + 0.21, 1.78, 0.16, desc, 6.8, GRAY, False, PP_ALIGN.CENTER)
add_shape(8.86, 1.61, 3.66, 0.37, GREEN_LIGHT, GREEN, MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(8.86, 1.64, 3.66, 0.16, "支撑技术：AI大模型自主规划 · 企业级多模态电力大模型", 7.8, GREEN_DARK, True, PP_ALIGN.CENTER)

# Middle left explanatory block
add_shape(0.45, 2.33, 3.72, 2.24, WHITE, BLUE, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0)
add_text(0.68, 2.47, 3.22, 0.6, "“设备—厂站—系统”\n三级数字孪生", 15, BLUE_DARK, True)
add_text(0.68, 3.18, 3.15, 0.76, "构建从微观设备到宏观系统的全域数字化映射，打通跨单位、跨专业壁垒，支撑全网资源优化配置。", 9.2, GRAY)
for i, item in enumerate(["物理电网与数字系统同步", "业务流程与控制策略耦合", "状态可视、风险可知、决策可算"]):
    dot(0.72, 4.02 + i * 0.23, GREEN)
    add_text(0.88, 3.95 + i * 0.23, 2.95, 0.18, item, 7.7, GREEN_DARK, i == 2)

# Central three-level digital twin tower
center_x = 6.5
# top intelligence halo
add_shape(center_x - 0.48, 0.98, 0.96, 0.96, GREEN_LIGHT, GREEN, MSO_SHAPE.OVAL, 1.5)
add_shape(center_x - 0.27, 1.19, 0.54, 0.54, WHITE, GREEN_MID, MSO_SHAPE.OVAL)
add_text(center_x - 0.27, 1.33, 0.54, 0.18, "AI", 14, GREEN_DARK, True, PP_ALIGN.CENTER)
add_line(center_x, 1.93, center_x, 2.12, GREEN, 1.2)

levels = [
    (5.25, 2.12, 2.50, 0.62, GREEN_MID, "系统级数字孪生", "全网态势 · 系统优化"),
    (4.82, 2.82, 3.36, 0.75, CYAN, "厂站级数字孪生", "站内协同 · 运行推演"),
    (4.38, 3.66, 4.24, 0.88, BLUE, "设备级数字孪生", "状态感知 · 寿命预测 · 风险诊断"),
]
for x, y, w, h, fill, title, sub in levels:
    add_shape(x, y, w, h, fill, fill, MSO_SHAPE.TRAPEZOID, 1.0)
    add_text(x + 0.22, y + 0.11, w - 0.44, 0.24, title, 11, WHITE, True, PP_ALIGN.CENTER)
    add_text(x + 0.22, y + 0.36, w - 0.44, 0.18, sub, 7.1, WHITE, False, PP_ALIGN.CENTER)
add_text(4.94, 4.61, 3.12, 0.28, "三级孪生贯通 · 全寿命周期联动", 9.5, GREEN_DARK, True, PP_ALIGN.CENTER)

# Middle right ID and resource coding block
add_shape(8.83, 2.33, 4.05, 2.24, WHITE, GREEN, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0)
add_text(9.08, 2.48, 3.55, 0.35, "实物“ID”与电网资源编码", 14, GREEN_DARK, True, PP_ALIGN.CENTER)
add_text(9.08, 2.88, 3.55, 0.4, "全域统一 · 同源更新 · 一物一码", 9, BLUE_DARK, True, PP_ALIGN.CENTER)
# Editable flow diagram
add_shape(9.18, 3.46, 0.78, 0.55, GREEN_LIGHT, GREEN, MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(9.18, 3.59, 0.78, 0.18, "实物流", 8, GREEN_DARK, True, PP_ALIGN.CENTER)
add_shape(10.48, 3.34, 0.78, 0.78, BLUE_LIGHT, BLUE, MSO_SHAPE.HEXAGON)
add_text(10.48, 3.58, 0.78, 0.2, "统一ID", 8, BLUE_DARK, True, PP_ALIGN.CENTER)
add_shape(11.78, 3.46, 0.78, 0.55, GREEN_LIGHT, GREEN, MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(11.78, 3.59, 0.78, 0.18, "价值流", 8, GREEN_DARK, True, PP_ALIGN.CENTER)
add_line(9.96, 3.74, 10.48, 3.74, GREEN, 1.5)
add_line(11.26, 3.74, 11.78, 3.74, GREEN, 1.5)
add_shape(10.48, 4.23, 0.78, 0.28, BLUE_LIGHT, BLUE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(10.48, 4.26, 0.78, 0.16, "信息流", 7.5, BLUE_DARK, True, PP_ALIGN.CENTER)

# Bottom ecosystem band
add_shape(0.45, 4.91, 12.43, 1.85, WHITE, GREEN_MID, MSO_SHAPE.ROUNDED_RECTANGLE, 1.0)
add_text(5.02, 5.0, 3.3, 0.28, "核心底座与融合生态", 12, GREEN_DARK, True, PP_ALIGN.CENTER)
add_line(0.72, 5.37, 12.62, 5.37, GRAY_LIGHT, 0.8)

bottom = [
    (0.72, 5.55, 3.35, "“技术—经济—环境”三重底线", "统筹安全可靠、投资效益与绿色低碳，追求全寿命周期综合最优。", ORANGE),
    (4.31, 5.55, 4.55, "AI原生与“全网一张图”底座", "同步设计物理电网、智能模型、业务流程与控制策略，实现实时耦合与深度融合。", GREEN),
    (9.10, 5.55, 3.48, "三流高度合一", "实物流、价值流、信息流贯通，形成穿透式闭环资产管理生态。", BLUE),
]
for x, y, w, title, desc, accent in bottom:
    add_shape(x, y, w, 0.96, RGBColor(250, 253, 252), accent, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(x + 0.15, y + 0.16, 0.43, 0.43, GREEN_LIGHT if accent != ORANGE else RGBColor(252, 238, 218), accent, MSO_SHAPE.OVAL)
    symbol = "底线" if accent == ORANGE else ("AI" if accent == GREEN else "三流")
    add_text(x + 0.15, y + 0.27, 0.43, 0.16, symbol, 6.8, accent, True, PP_ALIGN.CENTER)
    add_text(x + 0.7, y + 0.11, w - 0.82, 0.28, title, 10, accent, True, PP_ALIGN.CENTER)
    add_text(x + 0.7, y + 0.41, w - 0.82, 0.4, desc, 7.1, GRAY, False, PP_ALIGN.CENTER)

# Final outcome ribbon
add_shape(1.55, 6.94, 10.23, 0.34, GREEN, GREEN, MSO_SHAPE.ROUNDED_RECTANGLE)
add_text(1.55, 6.99, 10.23, 0.2, "最终实现：电网资产底盘绝对清晰 · 全寿命周期综合最优 · 超级智能体自主进化", 10, WHITE, True, PP_ALIGN.CENTER)

prs.save(OUT)
print(OUT)
